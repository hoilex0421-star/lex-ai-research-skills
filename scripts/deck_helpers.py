# -*- coding: utf-8 -*-
"""
deck_helpers.py — reusable python-pptx helpers for the Huawei-style
"market / insight 论点页" (16:9, 13.333 x 7.5 in).

Why this exists: every page of this template repeats the same 5 structural
parts (title band / insight band / left quant chart / right same-field cards /
multi-source footnote). Rebuilding these from scratch each time wastes effort
and invites inconsistency. Import this module and assemble pages from the parts.

Critical gotchas baked in (learned the hard way):
- Chinese text needs the 微软雅黑 typeface set on latin/ea/cs or it renders in a
  fallback font. `_setfont` handles this.
- For diagonal LINE charts use seg() (STRAIGHT connector, type 1). The ELBOW
  connector (type 2) renders a diagonal as ugly right-angle stairs.
- Render every page to PNG and eyeball it for label overlaps before shipping.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU = 914400
def IN(v): return Emu(int(v * EMU))
def C(h):  return RGBColor.from_string(h)

# --- palette: red is the ONLY accent colour; everything else is neutral ---
RED   = C('C7000B')   # primary brand red / 重点 / 工业主轴
RED2  = C('D85C5C')   # secondary red (商业 / 第二序列)
RED3  = C('EAB3B5')   # light red (家庭 / 第三序列)
DARK  = C('1A1A1A'); MID = C('595959'); MUTE = C('8C8C8C'); WHITE = C('FFFFFF')
GRID  = C('D9D9D9'); GF = C('F6F6F6'); PINK = C('FBECED'); GRAYH = C('6E6E6E'); LINEGRY = C('E2E2E2')
# --- secondary CATEGORY-chip palette (narrative pages only) ---
# Red stays the dominant accent + the ONLY colour used for conclusions/重点.
# Blue/green/navy are MUTED category tags (场景生成/场景理解, 输入/输出, 显式/隐式) — never for emphasis.
BLUE  = C('2E5C9E')   # category chip A
GREEN = C('4F8A3D')   # category chip B
NAVY  = C('1F3864')   # category chip C / axis fill
BLUEF = C('EAF1FA')   # very light blue panel fill
GREENF= C('EEF5E9')   # very light green panel fill
FONT  = "微软雅黑"

def newdeck():
    prs = Presentation(); prs.slide_width = IN(13.333); prs.slide_height = IN(7.5); return prs

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(s.placeholders): ph._element.getparent().remove(ph._element)
    return s

def _setfont(r, size, color, bold):
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    rPr = r._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        e = rPr.find(qn(tag))
        if e is None: e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', FONT)

def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=None, wrap=True, sb=0.0):
    """paras = list of paragraphs; each paragraph = list of (text, size, color, bold) runs."""
    tb = s.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h)); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor; tf.auto_size = MSO_AUTO_SIZE.NONE
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'): setattr(tf, m, 0)
    for k, pr in enumerate(paras):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align
        if ls: p.line_spacing = ls
        if sb: p.space_before = Pt(sb)
        for t, sz, col, b in pr:
            r = p.add_run(); r.text = t; _setfont(r, sz, col, b)
    return tb

def rect(s, x, y, w, h, fill, line=None, lw=0.75, rounded=False, rad=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             IN(x), IN(y), IN(w), IN(h))
    if rounded:
        try: shp.adjustments[0] = rad
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def hline(s, x, y, w, color=GRID, wpt=1.0, dash=None):
    ln = s.shapes.add_connector(2, IN(x), IN(y), IN(x + w), IN(y))
    ln.line.color.rgb = color; ln.line.width = Pt(wpt); ln.shadow.inherit = False
    if dash:
        d = ln.line._get_or_add_ln(); d.append(d.makeelement(qn('a:prstDash'), {'val': dash}))
    return ln

def seg(s, x1, y1, x2, y2, color=RED, wpt=2.2):
    """STRAIGHT line segment (connector type 1). Use for diagonal trend lines."""
    ln = s.shapes.add_connector(1, IN(x1), IN(y1), IN(x2), IN(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(wpt); ln.shadow.inherit = False
    return ln

def dot(s, cx, cy, d, color):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(cx - d / 2), IN(cy - d / 2), IN(d), IN(d))
    o.fill.solid(); o.fill.fore_color.rgb = color; o.line.fill.background(); o.shadow.inherit = False
    return o

# ---------------- the 5 structural parts ----------------

def title_band(s, red_part, black_part):
    """Part 1: 论点型标题 — red key clause + black remainder, with a quantified payload."""
    text(s, 0.5, 0.30, 12.4, 0.60, [[(red_part, 20, RED, True), (black_part, 20, DARK, True)]],
         anchor=MSO_ANCHOR.MIDDLE, ls=1.04)
    hline(s, 0.5, 0.96, 12.33, RED, 2.2)

def insight_band(s, tag, bullets, y=1.06, h=0.96):
    """Part 2: 「洞察启示」 — red tag + 2 conclusion-first bullets.
    tag = list of short lines for the red chip; bullets = list of paragraph run-lists."""
    rect(s, 0.5, y, 1.18, h, RED, rounded=True, rad=0.12)
    text(s, 0.5, y, 1.18, h, [[(t, 11.5, WHITE, True)] for t in tag],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.05)
    rect(s, 1.78, y, 11.05, h, GF, line=LINEGRY, lw=0.75, rounded=True, rad=0.06)
    paras = [[("▪ ", 9, RED, True)] + b for b in bullets]
    text(s, 1.98, y + 0.05, 10.7, h - 0.1, paras, anchor=MSO_ANCHOR.MIDDLE, ls=1.16, sb=3)

def card(s, x, y, w, h, name, note, chip, rows, strip_w=1.66):
    """Part 4: one same-field comparison card.
    name = card title; note = run-lists for the red strip (规模/关键数);
    chip = small colour dot tying to the chart legend; rows = list of (label, value).
    ALL cards on a page MUST share the same row labels so columns read as a comparison."""
    rect(s, x, y, w, h, WHITE, line=LINEGRY, lw=1.0, rounded=True, rad=0.05)
    rect(s, x + 0.04, y + 0.05, strip_w, h - 0.10, RED, rounded=True, rad=0.06)
    dot(s, x + 0.20, y + 0.27, 0.12, chip)
    text(s, x + 0.32, y + 0.17, strip_w - 0.30, 0.24, [[(name, 9.8, WHITE, True)]])
    text(s, x + 0.14, y + 0.40, strip_w - 0.18, h - 0.46, note, anchor=MSO_ANCHOR.TOP, ls=1.18)
    cx = x + strip_w + 0.14; cw = w - strip_w - 0.24
    paras = [[(lab + "  ", 7.2, RED, True), (val, 7.2, DARK, False)] for lab, val in rows]
    text(s, cx, y + 0.10, cw, h - 0.16, paras, anchor=MSO_ANCHOR.MIDDLE, ls=1.16, sb=2.5)

def footer(s, src, brand="市场/洞察论点页"):
    hline(s, 0.5, 7.06, 12.33, LINEGRY, 0.75)
    text(s, 0.5, 7.12, 9.5, 0.3, [[(src, 7, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 10.0, 7.12, 2.83, 0.3, [[(brand, 7.5, RED, True)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def stacked_bar(s, region, bars, series_colors, ymax, title=None):
    """Part 3 (one option): vertical stacked-bar chart for the left half.
    region = (x_left, y_top, x_right, baseline_y) in inches.
    bars = list of dicts: {x, label_top, label_bottom, vals:[...], total, top_is_red}
    series_colors = [c0, c1, ...] for the stacked segments (bottom->top).
    Caller positions bar x's; this draws segments, in-bar value labels, top totals."""
    xl, yt, xr, yb = region
    scale = (yb - yt) / float(ymax)
    hline(s, xl, yb, xr - xl, GRID, 1.0)
    for b in bars:
        x = b['x']; bw = b.get('bw', 0.82); cb = 0
        for vi, v in enumerate(b['vals']):
            if v <= 0: continue
            ty = yb - (cb + v) * scale; hh = v * scale
            rect(s, x, ty, bw, hh, series_colors[vi])
            if hh > 0.24:
                text(s, x, ty, bw, hh, [[(str(v), 7.6, WHITE, True)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cb += v
        topy = yb - b['total'] * scale
        col = RED if b.get('top_is_red') else GRAYH
        text(s, x - 0.1, topy - 0.26, bw + 0.2, 0.22, [[(str(b['total']), 10, col, True)]], align=PP_ALIGN.CENTER)
        if b.get('label_bottom'):
            text(s, x - 0.1, yb + 0.03, bw + 0.2, 0.2, [[(b['label_bottom'], 8, col, True)]], align=PP_ALIGN.CENTER)

# ============================================================================
# Technical-narrative helpers (timeline / framework-flow / tech-path / table /
# case pages). These model the second page family this skill covers. Red is
# still the dominant accent and the ONLY colour for conclusions; BLUE/GREEN/NAVY
# are muted CATEGORY chips. Every page still ends with a conclusion_band.
# ============================================================================

def _est_w(t, size):
    """Rough rendered width (inches) of a string at `size` pt for 微软雅黑."""
    cjk = sum(1 for ch in t if ord(ch) > 0x2E80)
    return (cjk * size * 1.02 + (len(t) - cjk) * size * 0.56) / 72.0

def tag_chip(s, x, y, label, fill=RED, txt=WHITE, size=8.5, w=None, h=0.26, rounded=True, pad=0.11):
    """Small category/label chip (场景生成 / 输入 / 显式表征 / 算力). Returns end-x.
    Auto-sizes width to the label unless `w` is given. Use fill=None for an
    outline-only chip (txt becomes the border+text colour)."""
    if w is None: w = _est_w(label, size) + 2 * pad
    if fill is None:
        rect(s, x, y, w, h, WHITE, line=txt, lw=1.0, rounded=rounded, rad=0.18)
        tc = txt
    else:
        rect(s, x, y, w, h, fill, rounded=rounded, rad=0.18); tc = txt
    text(s, x, y, w, h, [[(label, size, tc, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return x + w

def legend(s, x, y, items, size=8, gap=0.22, d=0.12):
    """Inline legend row: items = [(label, color), ...]. Returns end-x."""
    cx = x
    for lab, co in items:
        dot(s, cx + d / 2, y + 0.11, d, co)
        text(s, cx + d + 0.05, y, _est_w(lab, size) + 0.1, 0.22, [[(lab, size, DARK, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
        cx += d + 0.05 + _est_w(lab, size) + gap
    return cx

def section_header(s, x, y, w, label, color=RED, size=11, bar=False, barfill=PINK, h=0.34):
    """A colored bold sub-claim header introducing a panel/column.
    bar=True puts it on a light fill strip (like the template's 显式/隐式 headers)."""
    if bar: rect(s, x, y, w, h, barfill, rounded=True, rad=0.10)
    text(s, x + (0.12 if bar else 0), y, w - (0.24 if bar else 0), h,
         [[(label, size, color, True)]], anchor=MSO_ANCHOR.MIDDLE)
    return y + h

def stage_box(s, x, y, w, h, title, body=None, fill=WHITE, line=GRID,
              title_color=DARK, accent=None, tsize=9.5, bsize=8):
    """A process-stage box for framework/flow pages. `accent` draws a thin top
    strip in that colour (use to colour-code a stage group)."""
    rect(s, x, y, w, h, fill, line=line, lw=1.0, rounded=True, rad=0.08)
    if accent is not None: rect(s, x, y, w, 0.07, accent, rounded=True, rad=0.4)
    if body:
        text(s, x + 0.08, y + 0.06, w - 0.16, 0.30, [[(title, tsize, title_color, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        text(s, x + 0.10, y + 0.36, w - 0.20, h - 0.42, [[(body, bsize, MID, False)]],
             anchor=MSO_ANCHOR.TOP, ls=1.14)
    else:
        text(s, x + 0.06, y, w - 0.12, h, [[(title, tsize, title_color, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.1)
    return shp_right(x, w)

def shp_right(x, w): return x + w

def arrow(s, x, y, w, h, color=RED, direction='right'):
    """Block arrow / chevron between stages. direction in right/left/up/down."""
    shape = {'right': MSO_SHAPE.RIGHT_ARROW, 'left': MSO_SHAPE.LEFT_ARROW,
             'up': MSO_SHAPE.UP_ARROW, 'down': MSO_SHAPE.DOWN_ARROW}[direction]
    a = s.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    return a

def time_axis(s, x, y, w, labels, color=NAVY, size=8.5, band_h=0.30):
    """Horizontal time axis with evenly spaced tick labels (学术界/产业界 timeline).
    Draws a thin coloured band with white tick labels; returns the list of tick
    x-CENTERS so the caller can drop event cards above (学术界) / below (产业界)."""
    rect(s, x, y, w, band_h, color, rounded=True, rad=0.5)
    n = len(labels); centers = []
    for i, lab in enumerate(labels):
        cx = x + w * (i + 0.5) / n
        text(s, cx - 0.6, y, 1.2, band_h, [[(lab, size, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        centers.append(cx)
        if i: hline(s, x + w * i / n, y, 0, color)  # noop guard
    return centers

def event_card(s, x, y, w, h, head_runs, body=None, accent=RED, img=False):
    """A timeline/grid event card: a coloured header line (date + 红字产品名) plus
    body text, with an optional image-placeholder slot on top.
    head_runs = a run-list for the header paragraph (mix dark + red emphasis)."""
    rect(s, x, y, w, h, WHITE, line=LINEGRY, lw=1.0, rounded=True, rad=0.06)
    iy = y + 0.06
    if img:
        ih = min(0.62, h * 0.42); image_ph(s, x + 0.06, iy, w - 0.12, ih); iy += ih + 0.05
    text(s, x + 0.08, iy, w - 0.16, 0.30, [head_runs], anchor=MSO_ANCHOR.TOP, ls=1.08)
    if body:
        text(s, x + 0.08, iy + 0.28, w - 0.16, h - (iy - y) - 0.30,
             [[(body, 7.2, MID, False)]], anchor=MSO_ANCHOR.TOP, ls=1.12)
    return x + w

def image_ph(s, x, y, w, h, caption=None, label="图/截图"):
    """Light placeholder box for a screenshot/figure (web images can't auto-embed —
    the user drops the real asset in later). Shows a faint label + optional caption."""
    rect(s, x, y, w, h, C('F2F2F2'), line=GRID, lw=0.75, rounded=True, rad=0.04)
    text(s, x, y, w, h, [[(label, 7, MUTE, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if caption:
        text(s, x, y + h + 0.01, w, 0.18, [[(caption, 6.5, MUTE, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return y + h

def spec_table(s, x, y, w, col_w, headers, rows, header_fill=RED, header_txt=WHITE,
               row_h=0.30, size=8, hsize=8.5, zebra=True):
    """Comparison/spec table with a coloured header row (template's 对比表).
    col_w = list of column widths (inches) summing to ~w. rows = list of row lists.
    Returns the y just below the table."""
    cols = [x]
    for cw in col_w: cols.append(cols[-1] + cw)
    # header
    rect(s, x, y, sum(col_w), row_h, header_fill, line=header_fill, lw=0.5)
    for j, hd in enumerate(headers):
        text(s, cols[j], y, col_w[j], row_h, [[(hd, hsize, header_txt, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # body
    cy = y + row_h
    for ri, row in enumerate(rows):
        rh = row_h
        fill = GF if (zebra and ri % 2) else WHITE
        rect(s, x, cy, sum(col_w), rh, fill, line=LINEGRY, lw=0.5)
        for j, cell in enumerate(row):
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            runs = cell if isinstance(cell, list) else [(cell, size, DARK, j == 0)]
            text(s, cols[j] + 0.06, cy, col_w[j] - 0.12, rh, [runs],
                 align=al, anchor=MSO_ANCHOR.MIDDLE, ls=1.08)
        cy += rh
    # column separators
    for cx in cols:
        hline(s, cx, y, 0, LINEGRY)  # placeholder; vertical via rect borders already
    return cy

def conclusion_band(s, runs, y=6.60, x=0.5, w=12.33, h=0.42, fill=PINK, tag="结论"):
    """Full-width bottom takeaway band (the template's recurring red bottom line).
    runs = a run-list for the takeaway sentence (lead clause in RED bold).
    Set tag=None to drop the chip and run the sentence full-width."""
    rect(s, x, y, w, h, fill, line=C('F2D9DC'), lw=0.75, rounded=True, rad=0.10)
    tx = x + 0.16
    if tag:
        cw = _est_w(tag, 9) + 0.22
        rect(s, x + 0.12, y + 0.07, cw, h - 0.14, RED, rounded=True, rad=0.18)
        text(s, x + 0.12, y + 0.07, cw, h - 0.14, [[(tag, 9, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx = x + 0.12 + cw + 0.14
    text(s, tx, y, x + w - tx - 0.16, h, [runs], anchor=MSO_ANCHOR.MIDDLE, ls=1.12)
    return y

def note(s, x, y, w, txt, size=6.8):
    """Small gray caption / 注 line."""
    text(s, x, y, w, 0.3, [[(txt, size, MUTE, False)]], anchor=MSO_ANCHOR.TOP, ls=1.12)
